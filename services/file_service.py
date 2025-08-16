# 一時ファイル処理サービス（RAG機能用）
import os
import tempfile
import magic
import io
from typing import Dict, Any, BinaryIO
from fastapi import UploadFile, HTTPException
from pathlib import Path
import logging

# ドキュメント処理用のインポート
from docx import Document as DocxDocument
import PyPDF2
import openpyxl
from PIL import Image
from pptx import Presentation
import csv

# PDF処理の複数ライブラリ対応
import pdfplumber
import fitz  # PyMuPDF
import pytesseract
import cv2
import numpy as np

logger = logging.getLogger(__name__)

class FileService:
    """一時ファイル処理とテキスト抽出（永続保存なし）"""
    
    def __init__(self):
        # 設定値
        self.max_file_size = int(os.getenv("MAX_FILE_SIZE", "52428800"))  # 50MB
        self.allowed_extensions = os.getenv("ALLOWED_FILE_EXTENSIONS", 
                                          "pdf,docx,pptx,xlsx,csv,txt,md,png,jpg,gif").split(",")
        
        # MIME型とファイル拡張子のマッピング
        self.mime_mapping = {
            "application/pdf": "pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
            "text/csv": "csv",
            "text/plain": "txt",
            "text/markdown": "md",
            "image/png": "png",
            "image/jpeg": "jpg",
            "image/gif": "gif"
        }
        
        # Tesseract OCR設定（Windowsの場合）
        tesseract_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            "tesseract"  # PATH環境変数に設定されている場合
        ]
        
        for path in tesseract_paths:
            try:
                pytesseract.pytesseract.tesseract_cmd = path
                # テスト実行してパスが有効か確認
                pytesseract.get_tesseract_version()
                logger.info(f"Tesseract OCRが見つかりました: {path}")
                break
            except Exception:
                continue
        else:
            logger.warning("Tesseract OCRが見つかりません。OCR機能は無効になります。")
    
    async def validate_file(self, file: UploadFile) -> Dict[str, Any]:
        """ファイルのバリデーション"""
        try:
            # ファイルサイズチェック
            file_content = await file.read()
            file_size = len(file_content)
            
            if file_size > self.max_file_size:
                return {
                    "valid": False, 
                    "error": f"ファイルサイズが制限を超えています（制限: {self.max_file_size // 1024 // 1024}MB）"
                }
            
            if file_size == 0:
                return {"valid": False, "error": "空のファイルはアップロードできません"}
            
            # ファイル拡張子チェック
            file_extension = Path(file.filename).suffix.lower().lstrip('.')
            if file_extension not in self.allowed_extensions:
                return {
                    "valid": False, 
                    "error": f"許可されていないファイル形式です（許可形式: {', '.join(self.allowed_extensions)}）"
                }
            
            # MIME型チェック
            mime_type = magic.from_buffer(file_content, mime=True)
            if mime_type not in self.mime_mapping:
                return {"valid": False, "error": "不正なファイル形式です"}
            
            # 拡張子とMIME型の整合性チェック
            expected_extension = self.mime_mapping[mime_type]
            if file_extension != expected_extension:
                return {"valid": False, "error": "ファイル拡張子とファイル内容が一致しません"}
            
            # ファイルをリセット
            await file.seek(0)
            
            return {
                "valid": True,
                "file_size": file_size,
                "mime_type": mime_type,
                "extension": file_extension
            }
            
        except Exception as e:
            logger.error(f"ファイル検証エラー: {e}")
            return {"valid": False, "error": "ファイル検証中にエラーが発生しました"}
    
    async def process_uploaded_file_and_extract_text(self, file: UploadFile) -> Dict[str, Any]:
        """アップロードファイルを一時処理してテキスト抽出（元ファイルは削除）"""
        temp_file_path = None
        print(f"[DEBUG] ファイル処理開始: {file.filename}")
        try:
            # ファイルバリデーション
            validation_result = await self.validate_file(file)
            if not validation_result["valid"]:
                return {"success": False, "message": validation_result["error"]}
            
            # 一時ファイルに保存
            file_extension = validation_result["extension"]
            with tempfile.NamedTemporaryFile(
                suffix=f".{file_extension}", 
                delete=False
            ) as temp_file:
                temp_file_path = temp_file.name
                content = await file.read()
                temp_file.write(content)
            
            # テキスト抽出
            extracted_text = await self.extract_text_from_file(temp_file_path, file_extension)
            
            if not extracted_text.strip():
                return {
                    "success": False, 
                    "message": f"ファイルからテキストを抽出できませんでした（{file.filename}）"
                }
            
            logger.info(f"テキスト抽出成功: {file.filename} ({len(extracted_text)}文字)")
            
            return {
                "success": True,
                "extracted_text": extracted_text,
                "file_info": {
                    "original_filename": file.filename,
                    "file_size": validation_result["file_size"],
                    "file_type": file_extension,
                    "mime_type": validation_result["mime_type"],
                    "text_length": len(extracted_text)
                }
            }
            
        except Exception as e:
            logger.error(f"ファイル処理エラー: {e}")
            return {"success": False, "message": f"ファイル処理に失敗しました: {str(e)}"}
        
        finally:
            # 一時ファイルを必ず削除
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.remove(temp_file_path)
                    logger.debug(f"一時ファイル削除: {temp_file_path}")
                except Exception as cleanup_error:
                    logger.warning(f"一時ファイル削除エラー: {cleanup_error}")
    
    async def extract_text_from_file(self, file_path: str, file_type: str) -> str:
        """ファイルからテキストを抽出"""
        try:
            if file_type == "pdf":
                return await self._extract_from_pdf(file_path)
            elif file_type == "docx":
                return await self._extract_from_docx(file_path)
            elif file_type == "pptx":
                return await self._extract_from_pptx(file_path)  # PowerPoint対応
            elif file_type == "xlsx":
                return await self._extract_from_xlsx(file_path)
            elif file_type == "csv":
                return await self._extract_from_csv(file_path)
            elif file_type in ["txt", "md"]:
                return await self._extract_from_text(file_path)
            elif file_type in ["png", "jpg", "gif"]:
                return await self._extract_from_image(file_path)
            else:
                return ""
                
        except Exception as e:
            logger.error(f"テキスト抽出エラー ({file_type}): {e}")
            return ""
    
    async def _extract_from_pdf(self, file_path: str) -> str:
        """PDFからテキスト抽出（複数ライブラリ + OCR対応）"""
        logger.info(f"PDF分析開始 - ファイル: {file_path}")
        print(f"[DEBUG] PDF分析開始 - ファイル: {file_path}")
        
        # 複数の方法でテキスト抽出を試行
        methods = [
            ("PyPDF2", self._extract_with_pypdf2),
            ("pdfplumber", self._extract_with_pdfplumber),
            ("PyMuPDF", self._extract_with_pymupdf)
        ]
        
        # OCRが利用可能な場合のみ追加
        try:
            pytesseract.get_tesseract_version()
            methods.append(("OCR (Tesseract)", self._extract_with_ocr))
            logger.info("OCR機能が利用可能です")
        except Exception:
            logger.info("OCR機能は利用できません（テキストベースPDF抽出のみ）")
        
        for method_name, extract_func in methods:
            try:
                logger.info(f"PDF抽出方法: {method_name} を試行中...")
                print(f"[DEBUG] PDF抽出方法: {method_name} を試行中...")
                
                text = await extract_func(file_path)
                
                if text and len(text.strip()) > 50:  # 十分なテキストが抽出された
                    logger.info(f"PDF抽出成功: {method_name} で {len(text)}文字抽出")
                    return text
                elif text and len(text.strip()) > 0:
                    logger.warning(f"PDF抽出部分成功: {method_name} で {len(text)}文字抽出（少量）")
                    # 少量でも保存しておく（最終手段として使用）
                    last_resort_text = text
                else:
                    logger.warning(f"PDF抽出失敗: {method_name} でテキストが抽出されませんでした")
                    
            except Exception as e:
                logger.error(f"PDF抽出エラー ({method_name}): {type(e).__name__}: {e}")
                continue
        
        # すべての方法が失敗した場合、少量でも抽出されたテキストがあれば返す
        if 'last_resort_text' in locals():
            logger.info(f"最終手段として少量テキストを返却: {len(last_resort_text)}文字")
            return last_resort_text
        
        logger.error("PDF抽出完全失敗: すべての方法でテキスト抽出に失敗しました")
        return ""

    async def _extract_with_pypdf2(self, file_path: str) -> str:
        """PyPDF2を使用したPDF抽出"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            total_pages = len(pdf_reader.pages)
            is_encrypted = pdf_reader.is_encrypted
            
            logger.info(f"PyPDF2: ページ数={total_pages}, 暗号化={is_encrypted}")
            
            if is_encrypted:
                logger.warning("PyPDF2: 暗号化PDFです")
                return ""
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text += f"[ページ {page_num + 1}]\n{page_text}\n\n"
        
        return text
    
    async def _extract_with_pdfplumber(self, file_path: str) -> str:
        """pdfplumberを使用したPDF抽出（表やレイアウト対応）"""
        text = ""
        with pdfplumber.open(file_path) as pdf:
            logger.info(f"pdfplumber: ページ数={len(pdf.pages)}")
            
            for page_num, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text += f"[ページ {page_num + 1}]\n{page_text}\n\n"
                
                # 表も抽出を試行
                tables = page.extract_tables()
                if tables:
                    for table_num, table in enumerate(tables):
                        text += f"[ページ {page_num + 1} - 表 {table_num + 1}]\n"
                        for row in table:
                            if row:
                                text += " | ".join([cell or "" for cell in row]) + "\n"
                        text += "\n"
        
        return text
    
    async def _extract_with_pymupdf(self, file_path: str) -> str:
        """PyMuPDFを使用したPDF抽出"""
        text = ""
        doc = fitz.open(file_path)
        
        logger.info(f"PyMuPDF: ページ数={len(doc)}")
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text()
            
            if page_text and page_text.strip():
                text += f"[ページ {page_num + 1}]\n{page_text}\n\n"
        
        doc.close()
        return text
    
    async def _extract_with_ocr(self, file_path: str) -> str:
        """OCRを使用したPDF抽出（画像ベースPDF対応）"""
        try:
            # Tesseractが利用可能か確認
            pytesseract.get_tesseract_version()
        except Exception:
            logger.warning("OCR: Tesseract OCRが利用できません")
            return ""
        
        text = ""
        doc = fitz.open(file_path)
        
        logger.info(f"OCR: ページ数={len(doc)} のOCR処理を開始")
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # ページを画像に変換
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2倍解像度
            img_data = pix.tobytes("png")
            
            # PILで画像を読み込み
            image = Image.open(io.BytesIO(img_data))
            
            # OpenCVで前処理（コントラスト改善）
            img_array = np.array(image)
            if len(img_array.shape) == 3:
                img_array = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            
            # コントラスト改善
            img_array = cv2.threshold(img_array, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
            
            # OCRでテキスト抽出
            page_text = pytesseract.image_to_string(
                Image.fromarray(img_array),
                lang='jpn+eng',  # 日本語と英語
                config='--psm 6'  # 一様なテキストブロック
            )
            
            if page_text and page_text.strip():
                text += f"[ページ {page_num + 1} - OCR]\n{page_text}\n\n"
                logger.info(f"OCR: ページ {page_num + 1} で {len(page_text)}文字抽出")
            else:
                logger.warning(f"OCR: ページ {page_num + 1} でテキストが抽出されませんでした")
        
        doc.close()
        return text
    
    async def _extract_from_docx(self, file_path: str) -> str:
        """Wordドキュメントからテキストを抽出"""
        text = ""
        try:
            doc = DocxDocument(file_path)
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
        except Exception as e:
            logger.error(f"DOCX抽出エラー: {e}")
        return text
    
    async def _extract_from_pptx(self, file_path: str) -> str:
        """PowerPointからテキストを抽出"""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text += f"[スライド {slide_num + 1}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
        except Exception as e:
            logger.error(f"PPTX抽出エラー: {e}")
        return text
    
    async def _extract_from_xlsx(self, file_path: str) -> str:
        """Excelファイルからテキストを抽出"""
        text = ""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"[シート: {sheet_name}]\n"
                for row in sheet.iter_rows():
                    row_text = ""
                    for cell in row:
                        if cell.value:
                            row_text += str(cell.value) + "\t"
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
        except Exception as e:
            logger.error(f"XLSX抽出エラー: {e}")
        return text
    
    async def _extract_from_csv(self, file_path: str) -> str:
        """CSVファイルからテキストを抽出"""
        text = ""
        try:
            # UTF-8で試行
            with open(file_path, 'r', encoding='utf-8', newline='') as file:
                csv_reader = csv.reader(file)
                for row_num, row in enumerate(csv_reader):
                    if row_num == 0:
                        text += "[ヘッダー]\n"
                    text += "\t".join(row) + "\n"
        except UnicodeDecodeError:
            try:
                # Shift_JISで再試行
                with open(file_path, 'r', encoding='shift_jis', newline='') as file:
                    csv_reader = csv.reader(file)
                    for row_num, row in enumerate(csv_reader):
                        if row_num == 0:
                            text += "[ヘッダー]\n"
                        text += "\t".join(row) + "\n"
            except Exception as e:
                logger.error(f"CSV抽出エラー: {e}")
        except Exception as e:
            logger.error(f"CSV抽出エラー: {e}")
        return text
    
    async def _extract_from_text(self, file_path: str) -> str:
        """テキストファイルからテキストを抽出"""
        text = ""
        try:
            # UTF-8で試行
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
        except UnicodeDecodeError:
            try:
                # Shift_JISで再試行
                with open(file_path, 'r', encoding='shift_jis') as file:
                    text = file.read()
            except Exception as e:
                logger.error(f"テキスト抽出エラー: {e}")
        except Exception as e:
            logger.error(f"テキスト抽出エラー: {e}")
        return text
    
    async def _extract_from_image(self, file_path: str) -> str:
        """画像ファイルからメタデータを抽出（OCRは今後実装予定）"""
        text = ""
        try:
            with Image.open(file_path) as img:
                # 基本的なメタデータを抽出
                text += f"[画像情報]\n"
                text += f"サイズ: {img.size[0]}x{img.size[1]}px\n"
                text += f"モード: {img.mode}\n"
                text += f"形式: {img.format}\n"
                
                # EXIFデータがあれば抽出
                if hasattr(img, '_getexif') and img._getexif():
                    exif_data = img._getexif()
                    text += "EXIF情報:\n"
                    for tag, value in exif_data.items():
                        text += f"  {tag}: {value}\n"
                
                # 将来的にはOCRを実装予定
                text += "\n[注意: 画像内のテキスト認識は未実装です]\n"
                
        except Exception as e:
            logger.error(f"画像メタデータ抽出エラー: {e}")
        return text